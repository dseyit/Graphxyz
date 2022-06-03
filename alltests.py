from pptx import Presentation
from pptx.util import Inches
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

pptParameters={'reportName': 'Report name', 'slides': [{'slideTitle': 'Spectra', 'slideFigureCaptions': ['Figure caption 1', 'Figure caption 2', 'Figure caption 3', 'Figure caption 4'], 'slideFigureLocations': ['/Users/seyitliyev/Documents/Graphxyz/Reports/Slide 1/120K.png', '/Users/seyitliyev/Documents/Graphxyz/Reports/Slide 1/150K.png', '/Users/seyitliyev/Documents/Graphxyz/Reports/Slide 1/0.47mW.png', '/Users/seyitliyev/Documents/Graphxyz/Reports/Slide 1/0.47mWW.png'], 'slideFigureNames': ['120K', '150K', '0.47mW', '0.47mWW'],'slideFigureSizes': [[5,3],[5,3],[5,3],[5,3]]}, {'slideTitle': 'Dynamics', 'slideFigureCaptions': ['Figure caption 1','Figure caption 2'], 'slideFigureLocations': ['/Users/seyitliyev/Documents/Graphxyz/Reports/Slide 2/0.47mWW.png','/Users/seyitliyev/Documents/Graphxyz/Reports/Slide 2/78K.png'], 'slideFigureNames': ['0.47mWW'],'slideFigureSizes': [[5,3],[5,3]]}]}

prs=Presentation()

slideWidth=16
slideHeight=9

prs.slide_width = Inches(slideWidth)
prs.slide_height = Inches(slideHeight)

for i in range(len(pptParameters['slides'])):
    lyt = prs.slide_layouts[6] # choosing a slide layout
    slide = prs.slides.add_slide(lyt) # adding a slide
    slideParams = pptParameters['slides'][i]
    
    left = Inches(0)
    top = Inches(0)
    width = Inches(16)
    height = Inches(0.5)
    text_box=slide.shapes.add_textbox(left, top, width, height)
    tb=text_box.text_frame
    p = tb.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text=slideParams['slideTitle']
    font = run.font
    font.name = 'Calibri'
    font.size = Pt(25)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    #font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    for k in range(len(slideParams['slideFigureCaptions'])):
        figsize=slideParams['slideFigureCaptions']
        if len(slideParams['slideFigureCaptions'])==1 and k==0:
            width=Inches(9)
            left=Inches(0.75)
            top=Inches(1.5)
            
            cwidth=Inches(9)
            cheight = Inches(2)
            cleft=Inches(0.75)+Inches(slideParams['slideFigureSizes'][k][0])*cheight/Inches(slideParams['slideFigureSizes'][k][1])
            ctop=Inches(1.5)
        elif len(slideParams['slideFigureCaptions'])==2 and k==0:
            width=Inches(7.5)
            left=Inches(0.25)
            top=Inches(1)
            
            cwidth=Inches(7.5)
            cheight = Inches(1)
            cleft=Inches(0.25)
            ctop=Inches(1)+Inches(slideParams['slideFigureSizes'][k][1])*cwidth/Inches(slideParams['slideFigureSizes'][k][0])
        elif len(slideParams['slideFigureCaptions'])==2 and k==1:
            width=Inches(7.5)
            left=Inches(16-8)
            top=Inches(1)
            
            cwidth=Inches(7.5)
            cheight = Inches(1)
            cleft=Inches(16-8)
            ctop=Inches(1)+Inches(slideParams['slideFigureSizes'][k][1])*cwidth/Inches(slideParams['slideFigureSizes'][k][0])
        elif len(slideParams['slideFigureCaptions'])==4 and k==0:
            width=Inches(6)
            left=Inches(1)
            top=Inches(0.5)
            
            cwidth=Inches(6)
            cheight = Inches(1)
            cleft=Inches(1)
            ctop=Inches(0.5)+Inches(slideParams['slideFigureSizes'][k][1])*cwidth/Inches(slideParams['slideFigureSizes'][k][0])
        elif len(slideParams['slideFigureCaptions'])==4 and k==1:
            width=Inches(6)
            left=Inches(16-8.25)
            top=Inches(0.5)
            
            cwidth=Inches(6)
            cheight = Inches(1)
            cleft=Inches(16-8.25)
            ctop=Inches(0.5)+Inches(slideParams['slideFigureSizes'][k][1])*cwidth/Inches(slideParams['slideFigureSizes'][k][0])
        elif len(slideParams['slideFigureCaptions'])==4 and k==2:
            width=Inches(6)
            left=Inches(1)
            top=Inches(4.6)
            
            cwidth=Inches(6)
            cheight = Inches(1)
            cleft=Inches(1)
            ctop=Inches(4.6)+Inches(slideParams['slideFigureSizes'][k][1])*cwidth/Inches(slideParams['slideFigureSizes'][k][0])
        elif len(slideParams['slideFigureCaptions'])==4 and k==3:
            width=Inches(6)
            left=Inches(16-8.25)
            top=Inches(4.6)
            
            cwidth=Inches(6)
            cheight = Inches(1)
            cleft=Inches(16-8.25)
            ctop=Inches(4.6)+Inches(slideParams['slideFigureSizes'][k][1])*cwidth/Inches(slideParams['slideFigureSizes'][k][0])
        
        path=slideParams['slideFigureLocations'][k]
        img=slide.shapes.add_picture(path,left,top,width=width)
        
        text_box=slide.shapes.add_textbox(cleft, ctop, cwidth, cheight)
        tb=text_box.text_frame
        p = tb.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text=slideParams['slideFigureCaptions'][k]
        font = run.font
        font.name = 'Calibri'

prs.save(''.join([pptParameters['reportName'],".pptx"])) #saving file